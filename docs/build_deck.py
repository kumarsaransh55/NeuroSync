# Generates the NeuroSync jury pitch deck (.pptx) from our strategy/pitch docs.
# Run: python docs/build_deck.py   ->   docs/NeuroSync-Jury-Pitch.pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand palette ----
GREEN_DARK = RGBColor(0x0F, 0x3D, 0x2E)
GREEN_MID  = RGBColor(0x16, 0x65, 0x34)
ACCENT     = RGBColor(0x22, 0xC5, 0x5E)
INK        = RGBColor(0x11, 0x18, 0x27)
GRAY       = RGBColor(0x5B, 0x61, 0x6E)
LIGHT      = RGBColor(0xF6, 0xF9, 0xF8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARDBG     = RGBColor(0xEC, 0xFD, 0xF5)
RED        = RGBColor(0xB9, 0x1C, 0x1C)
FONT       = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
_num = 0

def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s

def box(s, l, t, w, h):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    return tf

def para(p, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, after=8, font=FONT):
    p.text = text
    p.alignment = align
    p.space_after = Pt(after)
    for run in p.runs:
        run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
        run.font.color.rgb = color; run.font.name = font
    return p

def first(tf):
    return tf.paragraphs[0]

def rect(s, l, t, w, h, color, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh

def accent(s):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), SH)
    b.fill.solid(); b.fill.fore_color.rgb = GREEN_DARK; b.line.fill.background(); b.shadow.inherit = False

def heading(s, kicker, title):
    accent(s)
    tf = box(s, 0.65, 0.5, 11.9, 1.3)
    if kicker:
        para(first(tf), kicker.upper(), 12, ACCENT, bold=True, after=3)
        para(tf.add_paragraph(), title, 30, GREEN_DARK, bold=True, after=0)
    else:
        para(first(tf), title, 30, GREEN_DARK, bold=True, after=0)

def bullets(s, items, l=0.75, t=2.05, w=11.9, h=4.4, size=18, gap=9):
    tf = box(s, l, t, w, h)
    for i, it in enumerate(items):
        p = first(tf) if i == 0 else tf.add_paragraph()
        para(p, "•   " + it, size, INK, after=gap)
    return tf

def source(s, text):
    tf = box(s, 0.65, 6.92, 11.4, 0.45)
    para(first(tf), "Source: " + text, 9.5, GRAY, italic=True, after=0)

def pagenum(s):
    global _num
    _num += 1
    tf = box(s, 12.5, 6.92, 0.6, 0.45)
    para(first(tf), str(_num), 10, GRAY, align=PP_ALIGN.RIGHT, after=0)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def style_cell(cell, text, size=13.5, color=INK, bold=False, bg=None, align=PP_ALIGN.LEFT):
    if bg is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.14)
    cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
    tf = cell.text_frame; tf.word_wrap = True
    para(first(tf), text, size, color, bold=bold, align=align, after=0)

def run(p, text, size, color, bold=False, italic=False, font=FONT, superscript=False, link=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    if superscript:
        r.font._rPr.set('baseline', '30000')  # superscript
    if link:
        r.hyperlink.address = link
    return r

# Canonical, numbered source list — used for the superscript fact-links AND the
# Sources slide, so the numbers match.
SOURCES = [
    (1, "Accenture — 'Getting to Equal: The Disability Inclusion Advantage' (2018)", "https://www.accenture.com/_acnmedia/pdf-89/accenture-disability-inclusion-research-report.pdf"),
    (2, "Accenture / Disability:IN / AAPD — 'The Disability Inclusion Imperative' (2023)", "https://newsroom.accenture.com/news/2023/companies-that-lead-in-disability-inclusion-outperform-peers-financially-reveals-new-research-from-accenture"),
    (3, "Government of India — Rights of Persons with Disabilities Act, 2016", "https://legislative.gov.in/sites/default/files/A2016-49.pdf"),
    (4, "nasscom — 'Neurodiversity and the Future of Work in India'", "https://community.nasscom.in/communities/neurodiversity/neurodiversity-and-future-work-india"),
    (5, "HRKatha — 'The 85% invisibility: Why India's neurodiverse talent stays hidden' (2024-25)", "https://www.hrkatha.com/features/research/the-85-invisibility-why-indias-neurodiverse-talent-stays-hidden/"),
    (6, "Frontiers in Psychology (2025) — universal support for neurodivergent employees", "https://www.frontiersin.org/journals/psychology/articles/10.3389/fpsyg.2025.1547877/full"),
    (7, "Branicki et al. — Human Resource Management, Wiley (2024)", "https://onlinelibrary.wiley.com/doi/10.1002/hrm.22243"),
    (8, "mydisabilityjobs.com — 'Neurodiversity in the Workplace | Statistics' (2025)", "https://mydisabilityjobs.com/statistics/neurodiversity-in-the-workplace/"),
    (9, "Disability:IN — evidence-based neuroinclusion framework", "https://disabilityin.org/articles-and-updates/disability-in-unveils-innovative-evidence-based-framework-to-advance-neuroinclusion-in-the-workforce"),
]
SRC_URL = {n: u for (n, _, u) in SOURCES}

# =====================================================================
# Slide 1 — Title
# =====================================================================
s = add_slide(GREEN_DARK)
# subtle accent block
rb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.9), SW, Inches(0.6))
rb.fill.solid(); rb.fill.fore_color.rgb = ACCENT; rb.line.fill.background(); rb.shadow.inherit = False
tf = box(s, 1.0, 2.3, 11.3, 2.6)
para(first(tf), "NeuroSync", 60, WHITE, bold=True, after=6)
para(tf.add_paragraph(), "Work, in sync with how you think.", 26, ACCENT, bold=True, after=14)
para(tf.add_paragraph(), "An inclusive work companion for neurodivergent employees (ADHD & dyslexia)", 18, WHITE, after=0)
tf2 = box(s, 1.0, 6.0, 11.3, 0.8)
para(first(tf2), "Presented by: [Your team]   ·   Accenture Inclusion Hackathon — Final Jury, 2026", 13, WHITE, after=0)
notes(s, "Hi, we're team NeuroSync. Before we show you anything we built, we want to tell you about one ordinary Monday morning. (Set a calm tone — one line, then move.)")
pagenum(s)

# =====================================================================
# Slide 2 — The human story
# =====================================================================
s = add_slide()
heading(s, "The human story", "Meet Aanya")
bullets(s, [
    "A talented analyst. She has ADHD and dyslexia.",
    "Monday, 9 AM: 38 unread emails, 4 chat threads, 6 tickets.",
    "One email hides three tasks and a Thursday deadline in the third paragraph.",
    "She re-reads it four times. Opens a ticket. A notification pulls her away.",
    "An hour later — she still hasn't started.",
], t=2.0, gap=11)
band = rect(s, 0.75, 5.75, 11.8, 1.0, CARDBG)
tf = box(s, 1.05, 5.9, 11.2, 0.8)
para(first(tf), "This isn't about Aanya's ability. Her workday was built for a brain that isn't hers.",
     18, GREEN_DARK, bold=True, after=0)
notes(s, "Tell Aanya's story slowly — this is the emotional hook. No tech yet. Land the last line: it's not about ability, it's about tools not designed for how she thinks.")
pagenum(s)

# =====================================================================
# Slide 3 — The problem
# =====================================================================
s = add_slide()
heading(s, "The problem", "A hidden 'cognitive tax'")
tf = box(s, 0.75, 1.95, 11.9, 1.7)
para(first(tf),
     "Turning workplace text into started, finished work costs employees with ADHD and dyslexia "
     "far more energy — not because of capability, but because mainstream tools assume one way of "
     "processing information.", 22, INK, bold=True, after=0)
bullets(s, [
    "Tasks, deadlines and priorities are buried inside long emails, chats and tickets.",
    "The result: missed deadlines, burnout and absence, and constant 'masking'.",
    "For the business: wasted talent and avoidable attrition.",
], t=3.95, gap=10)
source(s, "51% of neurodivergent employees have taken time off due to their condition — Neurodiversity in the Workplace statistics, mydisabilityjobs.com (2025).")
notes(s, "State the problem in one sentence (the bold block). Then the business impact. This is the 'problem statement' the jury asked for — say it before any solution.")
pagenum(s)

# =====================================================================
# Slide 4 — Who + how big (stats)
# =====================================================================
s = add_slide()
heading(s, "Who it's for — and why it matters", "A specific audience, a connected problem")
stats = [
    ("15–20%", "of people are neurodivergent (ADHD, autism, dyslexia, dyspraxia)", 8),
    ("~1 in 2", "people with ADHD also have dyslexia — so it's ONE connected problem, not two", 8),
    ("+28% / 2x / +30%", "revenue / net income / profit margin for disability-inclusion leaders", 1),
]
x = 0.75
for big, label, ref in stats:
    rect(s, x, 2.1, 3.78, 2.5, CARDBG)
    tf = box(s, x + 0.25, 2.35, 3.3, 2.1)
    p = first(tf)
    para(p, big, 30, GREEN_DARK, bold=True, after=8)
    run(p, " " + str(ref), 13, ACCENT, bold=True, superscript=True, link=SRC_URL[ref])  # clickable footnote -> source
    para(tf.add_paragraph(), label, 14, INK, after=0)
    x += 3.95
tf = box(s, 0.75, 4.95, 11.9, 1.4)
para(first(tf), "Our user: Accenture employees with ADHD and/or dyslexia.", 18, INK, bold=True, after=6)
para(tf.add_paragraph(),
     "Accenture's own research already proved inclusion pays. NeuroSync operationalises it for the "
     "neurodivergent talent already inside our walls.", 15, GRAY, italic=True, after=0)
source(s, "Accenture, 'Getting to Equal: The Disability Inclusion Advantage' (2018). Prevalence & ADHD/dyslexia comorbidity: mydisabilityjobs.com (2025).")
notes(s, "Lead the financial stat to an Accenture jury — it's Accenture's own research. The comorbidity stat defuses the 'you're doing two disabilities' critique: it's one connected problem.")
pagenum(s)

# =====================================================================
# Slide 5 — Solution + 3 USPs
# =====================================================================
s = add_slide()
heading(s, "Our solution", "NeuroSync — and what makes it different")
tf = box(s, 0.75, 1.85, 11.9, 0.7)
para(first(tf), "An inclusive work companion that turns the chaos of workplace communication into "
     "clarity and momentum — adapted to how each person's brain works.", 16, INK, after=0)
usps = [
    ("1  Adaptive output", "\"Same message, rendered for your brain.\" The same email becomes time-boxed steps (ADHD) or short, plain, dyslexia-friendly text — Copilot gives everyone the identical summary."),
    ("2  Understanding -> doing", "\"From 'I get it' to 'I started.'\" We attack task initiation and overwhelm, not just comprehension: micro-steps, focus mode, a one-tap reset."),
    ("3  Inclusion-first & enterprise-native", "\"Built for your brain AND your firewall.\" Calm, sensory-aware by default, running on enterprise data behind governance — not a public chatbot."),
]
x = 0.75
for title, body in usps:
    rect(s, x, 2.75, 3.78, 3.7, LIGHT, line=RGBColor(0xD9,0xE5,0xDF))
    tf = box(s, x + 0.25, 3.0, 3.3, 3.3)
    para(first(tf), title, 16, GREEN_DARK, bold=True, after=8)
    para(tf.add_paragraph(), body, 13, INK, after=0)
    x += 3.95
notes(s, "Three things the incumbents don't do together. This is the differentiation the jury demanded. Then bridge: 'instead of more slides, let us show you.'")
pagenum(s)

# =====================================================================
# Slide 6 — Demo
# =====================================================================
s = add_slide()
heading(s, "See it work", "Live demo — real AI, in seconds")
bullets(s, [
    "Paste Aanya's messy email — or pull it straight from Outlook (Microsoft Graph).",
    "NeuroSync returns a plain summary, the real action items, and the hidden Thursday deadline.",
    "Toggle Dyslexia mode — the same content, re-rendered for a dyslexic reader.",
    "Send it to the Planner — AI breaks it into ordered micro-steps with time estimates.",
    "Focus Mode strips everything but the next step. 'Feeling overwhelmed?' gives a calm reset.",
], t=2.0, gap=12)
band = rect(s, 0.75, 5.95, 11.8, 0.85, CARDBG)
tf = box(s, 1.05, 6.08, 11.2, 0.65)
para(first(tf), "Running live on real AI (Google Gemini today; Azure OpenAI in production).",
     15, GREEN_DARK, bold=True, after=0)
notes(s, "Switch to the running app and follow DEMO-SCRIPT.md. Keep the backup recording ready. End: 'same email — now understood, broken down, started, and she never had to mask or ask for help.'")
pagenum(s)

# =====================================================================
# Slide 7 — Why us vs Copilot/ChatGPT
# =====================================================================
s = add_slide()
heading(s, "Differentiation", "Why NeuroSync, not Copilot / ChatGPT")
rows = [
    ("Capability", "NeuroSync", "Copilot / ChatGPT"),
    ("Adapts output to your cognitive profile", "Yes", "No"),
    ("Helps you START, not just understand", "Yes", "No"),
    ("In-the-moment regulation (overwhelm reset)", "Yes", "No"),
    ("Runs on enterprise data + governance", "Yes", "Partial"),
]
t = s.shapes.add_table(len(rows), 3, Inches(0.75), Inches(2.1), Inches(11.8), Inches(3.6)).table
t.columns[0].width = Inches(6.2); t.columns[1].width = Inches(2.8); t.columns[2].width = Inches(2.8)
for c in range(3):
    style_cell(t.cell(0, c), rows[0][c], size=15, color=WHITE, bold=True, bg=GREEN_DARK,
               align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
for r in range(1, len(rows)):
    rowbg = WHITE if r % 2 else LIGHT
    style_cell(t.cell(r, 0), rows[r][0], size=14, bg=rowbg)
    style_cell(t.cell(r, 1), rows[r][1], size=14, bold=True, color=GREEN_MID, bg=rowbg, align=PP_ALIGN.CENTER)
    col = GRAY if rows[r][2] == "No" else RGBColor(0x9A,0x6A,0x00)
    style_cell(t.cell(r, 2), rows[r][2], size=14, color=col, bg=rowbg, align=PP_ALIGN.CENTER)
notes(s, "Generic AI summarises. NeuroSync adapts to the person, helps them act, supports them in the moment, and respects enterprise data. That combination is the moat.")
pagenum(s)

# =====================================================================
# Slide 8 — Responsible & inclusive by design
# =====================================================================
s = add_slide()
heading(s, "Ethics & inclusion", "Responsible and inclusive by design")
bullets(s, [
    "No diagnosis, no labels — support is an opt-in user choice, never an assumption.",
    "Human-in-the-loop — the AI suggests; the person edits and approves. Nothing auto-acts.",
    "Sourced nudges — wellbeing prompts are based on recognised, published techniques.",
    "Privacy by design — enterprise data stays within enterprise governance.",
    "Inclusive by default — calm UI, dyslexia font, high contrast, reduced motion for everyone.",
], t=2.05, gap=13)
notes(s, "We're careful because this community is sensitive. This slide directly answers 'ethical/responsible AI' — a scored criterion. Keep it crisp.")
pagenum(s)

# =====================================================================
# Slide 9 — Feasibility, security & scale
# =====================================================================
s = add_slide()
heading(s, "Feasibility & security", "Real today — and enterprise-ready")
bullets(s, [
    "Working MVP: React frontend + ASP.NET Core API + SQL on Azure, with real AI (Google Gemini).",
    "Real Outlook integration via Microsoft Graph (email -> AI-broken-down task).",
    "Security: in production the AI runs on Azure OpenAI inside our tenant — client data never leaves Accenture's compliance boundary; self-hosted open models for the most sensitive data.",
    "The AI sits behind an interface, so switching providers is a config change, not a rewrite.",
    "Scales as a lightweight enterprise layer; the adaptive engine extends to more needs.",
], t=2.05, gap=12)
source(s, "Architecture & security detail: repo docs/SECURITY-AND-AI.md and docs/OUTLOOK-INTEGRATION.md.")
notes(s, "Pre-empt the data-security question: 'for the demo we use Gemini; in production this runs on Azure OpenAI in our own tenant.' Covers privacy + feasibility + responsible AI in one breath.")
pagenum(s)

# =====================================================================
# Slide 10 — "You asked, we listened" (defense)
# =====================================================================
s = add_slide()
heading(s, "Acting on your feedback", "How we sharpened it since the last review")
rows = [
    ("Your feedback", "What we changed"),
    ("No clear problem statement", "Problem-first story (Aanya) + a one-line problem and business impact"),
    ("Target audience too vague", "Accenture employees with ADHD/dyslexia — a named persona"),
    ("'Productivity' implies they're less productive", "Reframed as an inclusive companion that removes friction"),
    ("Overlaps with Copilot / ChatGPT", "Three concrete USPs the incumbents don't combine"),
    ("Needs validation", "User interviews + credible, cited research"),
    ("Name not justified", "Neuro + Sync = sync work to how your brain works"),
]
t = s.shapes.add_table(len(rows), 2, Inches(0.75), Inches(2.0), Inches(11.8), Inches(4.4)).table
t.columns[0].width = Inches(5.0); t.columns[1].width = Inches(6.8)
for c in range(2):
    style_cell(t.cell(0, c), rows[0][c], size=15, color=WHITE, bold=True, bg=GREEN_DARK)
for r in range(1, len(rows)):
    rowbg = WHITE if r % 2 else LIGHT
    style_cell(t.cell(r, 0), rows[r][0], size=13, bg=rowbg, color=GRAY)
    style_cell(t.cell(r, 1), rows[r][1], size=13, bg=rowbg, bold=True)
notes(s, "This slide is for THIS jury — it shows we listened. Optional to present in full; great to have for Q&A. Move through it quickly with confidence.")
pagenum(s)

# =====================================================================
# Slide 11 — Roadmap
# =====================================================================
s = add_slide()
heading(s, "What's next", "From prototype to pilot")
cols = [
    ("NOW — working", ["AI task breakdown", "Plain-language summarizer", "Outlook email -> task", "Focus mode + overwhelm reset"]),
    ("NEXT — 1-3 months", ["Validated pilot with a neurodiversity ERG", "Azure OpenAI in-tenant", "Teams & Jira connectors"]),
    ("LATER — scale", ["More language profiles", "Low-vision / screen-reader modes", "Manager insights (privacy-safe)"]),
]
x = 0.75
for title, items in cols:
    rect(s, x, 2.1, 3.78, 4.0, LIGHT, line=RGBColor(0xD9,0xE5,0xDF))
    tf = box(s, x + 0.25, 2.3, 3.35, 3.7)
    para(first(tf), title, 15, GREEN_DARK, bold=True, after=10)
    for it in items:
        para(tf.add_paragraph(), "•  " + it, 13, INK, after=7)
    x += 3.95
notes(s, "What you saw is real and on Azure. The near-term ask is a validated pilot. Connectors and more accessibility profiles follow.")
pagenum(s)

# =====================================================================
# Slide 12 — Close + ask
# =====================================================================
s = add_slide(GREEN_DARK)
tf = box(s, 1.0, 2.3, 11.3, 3.0)
para(first(tf), "Aanya doesn't need to be fixed. Her workday does.", 32, WHITE, bold=True, after=18)
para(tf.add_paragraph(), "NeuroSync levels the playing field — quietly, respectfully, and at scale.", 20, ACCENT, after=22)
para(tf.add_paragraph(), "Our ask: help us pilot NeuroSync with a neurodiversity employee resource group.", 18, WHITE, after=0)
tf2 = box(s, 1.0, 6.1, 11.3, 0.8)
para(first(tf2), "NeuroSync  ·  Work, in sync with how you think.", 15, WHITE, bold=True, after=0)
notes(s, "Land the mission line. Make a concrete ask (a pilot). Thank them and invite questions — see QA-PREP.md.")
pagenum(s)

# =====================================================================
# Slide 13 — Sources & references
# =====================================================================
s = add_slide()
heading(s, "Credibility", "Sources & references")
tf = box(s, 0.75, 1.95, 11.9, 4.7)
for i, (num, label, url) in enumerate(SOURCES):
    p = first(tf) if i == 0 else tf.add_paragraph()
    para(p, f"{num}.  {label} — ", 12, INK, after=6)
    run(p, url, 10.5, ACCENT, link=url)  # clickable
tf2 = box(s, 0.75, 6.55, 11.9, 0.5)
para(first(tf2), "Tier-1 sources (Accenture / industry bodies / peer-reviewed) are prioritised; verify exact figures against primary sources before the final.",
     10.5, GRAY, italic=True, after=0)
notes(s, "Keep this slide available for Q&A. If a juror challenges a number, point to the primary source.")
pagenum(s)

out = __file__.rsplit("\\", 1)[0] + "\\NeuroSync-Jury-Pitch.pptx"
prs.save(out)
print("SAVED", out, "slides=", len(prs.slides._sldIdLst))
